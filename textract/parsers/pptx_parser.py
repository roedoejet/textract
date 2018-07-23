import pptx

from .utils import BaseParser


class Parser(BaseParser):
    """Extract text from pptx file using python-pptx
    """
    def __init__(self, **kwargs):
        print(kwargs)

    def extract(self, filename, **kwargs):
        if "language" in kwargs and kwargs['language']:
            cors = processCors(kwargs["language"])
        else:
            print("Extracting unconverted text because no conversion language was specified.")
        presentation = pptx.Presentation(filename)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return '\n\n'.join(text_runs)



        with open(filename) as stream:
            text = stream.read()
            if "language" in kwargs and kwargs['language']:
                text = cors.apply_rules(text)
            return text